# -*- coding: utf-8 -*-
"""
Created on Tue May 16 11:22:16 2023

@author: fecia
"""

import numpy as np
import pydicom
import cv2
import os
import imageio
import glob
import pptx
import tkinter as tk
from tkinter import ttk
from tkinter import filedialog
import sys
import imageio_ffmpeg

# Configuration of the FFmpeg executable
if getattr(sys, 'frozen', False):
    dir_path = os.path.dirname(sys.executable)
else:
    dir_path = os.path.dirname(os.path.abspath(__file__))

ffmpeg_exe = os.path.join(dir_path, 'ffmpeg.exe')
imageio_ffmpeg.get_ffmpeg_exe = lambda: ffmpeg_exe

# Verification of the FFmpeg path
print("FFmpeg executable used by imageio:", imageio_ffmpeg.get_ffmpeg_exe())

def process_dicom_files(folder_path):
    global prs
    # Create the presentation object
    prs = pptx.Presentation()
    num_diapo = 0
    
    # Search for all files in the specified folder and its subfolders
    dcm_files = [f for f in glob.glob(folder_path + "/**/*", recursive=True) if os.path.isfile(f)]
    
    # Sort the dcm_files list based on SeriesNumber
    sorted_dcm_files = []
    for idx, dcm_file in enumerate(dcm_files):
        try:
            dcm_data = pydicom.dcmread(dcm_file, force=True)
            series_number = getattr(dcm_data, 'SeriesNumber', None)
            if series_number is not None:
                sorted_dcm_files.append((series_number, dcm_file))
        except Exception as e:
            print(f"Error reading {dcm_file}: {e}")
            continue
    
    sorted_dcm_files.sort()
    sorted_dcm_files = [x[1] for x in sorted_dcm_files]
    
    total_files = len(sorted_dcm_files)
    
    # Process each file
    for i, dcm_file in enumerate(sorted_dcm_files):
        # Update progress
        progress = (i + 1) / total_files * 100
        progress_var.set(progress)
        message_label.config(text=f"Processing file {i+1} of {total_files}: {os.path.basename(dcm_file)}")
        root.update()
        
        # Open the DICOM file
        try:
            dcm_data = pydicom.dcmread(dcm_file, force=True)
        except Exception as e:
            print(f"Error reading {dcm_file}: {e}")
            continue
        
        modality = getattr(dcm_data, 'Modality', '')
        if modality != "XA":
            # If the modality is not "XA", skip this file
            print(f"{dcm_file} modality is {modality}, skipping.")
            continue
        
        # Check if the DICOM file contains the necessary pixel data
        if "PixelData" not in dcm_data:
            print(f"{dcm_file} does not contain pixel data, skipping.")
            continue
        
        # Process the DICOM data
        pixel_data = dcm_data.pixel_array
        
        # If the pixel data is 2D, it's an image
        if pixel_data.ndim == 2:
            # Save the image
            image_filename = f"image_{i+1}.png"
            cv2.imwrite(image_filename, pixel_data)
            # Add a slide to the presentation
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            # Add the image to the slide
            picture = slide.shapes.add_picture(
                image_filename, pptx.util.Cm(1), pptx.util.Cm(1), width=pptx.util.Cm(14), height=pptx.util.Cm(14)
            )
        # If the pixel data is 3D, it's a video
        elif pixel_data.ndim == 3:
            # Save video
            # Split the file name and extension
            file_name, file_ext = os.path.splitext(dcm_file)
            # Add the prefix and iteration number to the file name
            new_file_name = f"video_{i+1}"
            # Paths for the video and poster image
            new_file_path = os.path.join(folder_path, new_file_name + ".mp4")
            poster_path = os.path.join(folder_path, new_file_name + ".png")
            # Save the video
            imageio.mimwrite(new_file_path, np.expand_dims(pixel_data, axis=-1))
            # Save poster image (middle frame)
            middle_frame = len(pixel_data) // 2
            cv2.imwrite(poster_path, pixel_data[middle_frame])
            # Create and save additional images
            for j in range(9):
                # Create image using the video data
                image = pixel_data[j * len(pixel_data) // 9]
                # Save the image to a file
                cv2.imwrite(f"image_{j+1}.png", image)
            # Add a slide to the presentation
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            # Add the video to the slide
            movie = slide.shapes.add_movie(
                new_file_path, pptx.util.Cm(13), pptx.util.Cm(3), width=pptx.util.Cm(11), height=pptx.util.Cm(11),
                mime_type='video/mp4', poster_frame_image=poster_path
            )
            # Add the images to the slide in a 3x3 grid
            for j in range(3):
                for k in range(3):
                    # Calculate the position of the image
                    x = pptx.util.Cm(0.5 + k * 4)
                    y = pptx.util.Cm(3 + j * 4)
                    # Add the image to the slide
                    picture = slide.shapes.add_picture(
                        f"image_{j*3+k+1}.png", x, y, width=pptx.util.Cm(4), height=pptx.util.Cm(4)
                    )
        else:
            print(f"Unknown pixel data format in {dcm_file}")
            continue
        
        num_diapo += 1
        print(f"Slide {num_diapo} added successfully")
        # Update GUI
        message_label.config(text=f"Slide {num_diapo} added successfully")
        root.update()
    
    # After processing is done
    message_label.config(text="Processing completed.")
    root.update()
    # Save the presentation
    save_presentation()

def save_presentation():
    global prs
    message_label.config(text="Saving presentation...")
    root.update()
    pres_name = filedialog.asksaveasfilename(defaultextension='.pptx', parent=root)
    if not pres_name:
        message_label.config(text="No file name provided. Presentation not saved.")
        start_button.config(state=tk.NORMAL)
        return
    prs.save(pres_name)
    message_label.config(text=f"Presentation saved as {pres_name}")
    start_button.config(state=tk.NORMAL)

def start_processing():
    # Disable the start button
    start_button.config(state=tk.DISABLED)
    # Open the "Open" dialog box and allow the user to select a folder
    folder_path = filedialog.askdirectory(parent=root)
    if not folder_path:
        message_label.config(text="No folder selected.")
        start_button.config(state=tk.NORMAL)
        return
    # Now start processing
    process_dicom_files(folder_path)

def main():
    global root, progress_var, message_label, start_button
    # Create the root window
    root = tk.Tk()
    root.title("DICOM to PowerPoint Converter")
    root.geometry("400x200")
    
    # Create a label to show messages
    message_label = tk.Label(root, text="Select a folder containing DICOM files.")
    message_label.pack(pady=10)
    
    # Create a progress bar
    progress_var = tk.DoubleVar()
    progress_bar = ttk.Progressbar(root, variable=progress_var, maximum=100)
    progress_bar.pack(fill=tk.X, padx=20, pady=10)
    
    # Create a button to start the processing
    start_button = tk.Button(root, text="Select Folder and Start", command=start_processing)
    start_button.pack(pady=10)
    
    # Run the main loop
    root.mainloop()

if __name__ == "__main__":
    main()
